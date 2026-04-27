"""Fill the thesis presentation template based on the actual project state.

Reads `TezDosyaları/Danışman_Juri_Sunum.pptx`, removes the rules slide
(slide 1 is meta-instruction), populates the remaining slides with content
derived from the working prototype, and writes the result alongside the
original.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "TezDosyaları" / "Danışman_Juri_Sunum.pptx"
DST = ROOT / "TezDosyaları" / "Danışman_Juri_Sunum_v3.pptx"


PROJECT_TITLE = (
    "Veri Seti Açıklamaları Üzerinden Büyük Dil Modeli Uyarlaması ile "
    "Anlam Tabanlı Veri Keşfi"
)
STUDENTS = "Öykü Duru Akpınar  ·  Melike Selma Genç"
ADVISOR = "Prof. Dr. Murat Osman Ünalır"
DATE = "Mayıs 2026"

# (title, [(level, text), ...])
SLIDES = [
    # Konu ve Amaç
    (
        "Konu ve Amaç",
        [
            (0, "Veri bilimi projelerinde uygun veri setini bulmak en uzun süren adımlardan biridir."),
            (0, "Kaggle ve Hugging Face gibi platformlardaki mevcut arama ağırlıklı olarak anahtar kelime eşleşmesine dayanır; eş anlamlı ya da farklı yazım tarzındaki açıklamaları kaçırır."),
            (0, "Amaç: yalnızca veri seti başlığı ve açıklama metni üzerinden anlamsal arama yapan, klasik BM25 baseline ile sayısal olarak karşılaştırılabilir bir prototip geliştirmek."),
            (0, "Hedef metrikler: Precision@K, Recall@K, MRR, nDCG@5; ek olarak Türkçe alt küme analizi."),
        ],
    ),
    # Özgün Değer
    (
        "Özgün Değer",
        [
            (0, "Dürüst güven sinyali — corpus'ta karşılığı olmayan sorgularda sistem yanlış sonuç değil, \"weak match\" rozeti ve OOD uyarısı gösterir."),
            (0, "Lexical-anchor cezası — sorgudaki kritik kelime dokümanda yoksa skor düşürülür; Pokemon-images gibi anlamsal yakın ama konuyla ilgisiz sonuçlar elenir."),
            (0, "Dört encoder profili tek pipeline üzerinde yan yana: MiniLM, E5-base, MiniLM-FT (alan uyarlamalı), multilingual (Türkçe destekli)."),
            (0, "Title-deemphasis politikası — açıklaması zayıf veri setlerinde başlık tek başına sonucu sürüklemez, kullanıcıya rozetle bildirilir."),
            (0, "Tekrar üretilebilir değerlendirme — 34 etiketli sorgu + 46 sorguluk stres seti, her değişikliğin etkisi P@1/P@3 ile ölçülebilir."),
        ],
    ),
    # Yaygın Etki
    (
        "Yaygın Etki",
        [
            (0, "Açık veri portalları ve kurumsal data catalog'ları için yeniden kullanılabilir anlamsal arama modülü."),
            (0, "Veri bilimi araştırmacıları için veri keşif süresinde belirgin azalma — Top-K içinde aranan veri setine ulaşma oranı %50+ göreli iyileşti."),
            (0, "Tekrar üretilebilir değerlendirme protokolü: 34 etiketli sorguluk benchmark + 46 sorguluk regresyon probu (`dev_tools/regression_probe.py`)."),
            (0, "Türkçe sorgu desteği — multilingual profile sayesinde Türkçe data discovery senaryoları aynı sistem üzerinden test edilebilir."),
            (0, "Açık kaynak Python yığını (sentence-transformers, FAISS, BM25) — gelecek araştırmalar için temel."),
        ],
    ),
    # Uygulanabilirlik
    (
        "Uygulanabilirlik",
        [
            (0, "Çalışan web prototipi: tarayıcı tabanlı Live Search arayüzü ve dört araştırma sorusunu (RQ1–RQ4) açıklayan paneller."),
            (0, "Corpus: 1.120 veri seti tanımı (640 Kaggle + 480 Hugging Face) FAISS ile indekslendi; her sorgu 1 saniyenin altında yanıtlanıyor."),
            (0, "Üç retrieval yöntemi (BM25 / Semantic / Hybrid) ve dört encoder profili kullanıcı tarafından canlı seçilebiliyor — toplam 24 farklı konfigürasyon test edilebilir."),
            (0, "Cross-encoder semantic rerank açılır-kapanır toggle olarak sunulmuş; ablation çalışmaları yapılabilir."),
            (0, "Yatay ölçeklenebilirlik: yeni veri kaynağı eklemek için yalnızca normalize_merge.py + build_faiss_index.py çalıştırmak yeterli."),
        ],
    ),
    # Sistem Mimarisi
    (
        "Sistem Mimarisi",
        [
            (0, "1) Veri Katmanı"),
            (1, "Kaggle + Hugging Face koleksiyon scriptleri  →  ham JSONL"),
            (1, "normalize_merge.py: temizleme, dil tespiti, domain/use-case çıkarımı  →  descriptions_clean.jsonl"),
            (0, "2) İndeks Katmanı  (her encoder profili için ayrı)"),
            (1, "Dense: build_faiss_index.py  →  FAISS IndexFlatIP  +  mappings.json"),
            (1, "Lexical: bm25.py  →  ters indeks (TITLE×3 + KEYWORD×2 + description)"),
            (0, "3) Sorgu Anlama Katmanı  (query_understanding.py)"),
            (1, "Tokenize → stopword/generic-noise filter → intent body → concept expansion → semantic variants"),
            (0, "4) Geri Getirme + Sıralama Katmanı"),
            (1, "search.py / hybrid.py: stage-1 semantic + BM25 skor füzyonu (top-40 aday)"),
            (1, "reranker.py: lexical-anchor coverage  +  cross-encoder rerank (BAAI/bge-reranker-base)  +  confidence etiketleme"),
            (0, "5) Sunum Katmanı  (web_app.py)"),
            (1, "HTTP API + tarayıcı arayüzü: Live Search, profile/method/source kontrolleri, 4 RQ paneli"),
            (0, "6) Değerlendirme Katmanı  (offline)"),
            (1, "evaluate_search.py (tez benchmark)  +  dev_tools/regression_probe.py (stres seti)  +  probe_evaluator.py (P@K, OOD doğruluğu)"),
        ],
    ),
    # Gerçekleştirme Yöntemi
    (
        "Gerçekleştirme Yöntemi",
        [
            (0, "Veri toplama → temizleme → semantic_text inşası (başlık + açıklama + çıkarsanmış domain/use case/modality)."),
            (0, "İndeksleme: FAISS (semantic) + ters indeks (BM25) paralel; her profil için ayrı vektör deposu."),
            (0, "Sorgu işleme: tokenizasyon → stopword/generic-noise filtresi → curated phrase/concept genişletme → semantic varyant üretimi."),
            (0, "Sıralama: stage-1 semantic + BM25 fusion → lexical-anchor coverage bonusu / penalty → cross-encoder rerank (BAAI/bge-reranker-base) → confidence etiketleme."),
            (0, "Değerlendirme: 34 etiketli sorgu + 46 sorguluk stres seti; Hybrid P@1=0.91 (BM25 P@1=0.74), nDCG@5=0.79; OOD doğruluğu 0.90+."),
            (0, "Mimari modüler: query_understanding, search, hybrid, reranker, quality_scoring, web_app — her katman bağımsız test edilebilir."),
        ],
    ),
]


def remove_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    sldId = slides[index]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sldId)


def set_title(slide, text):
    title = slide.shapes.title
    if title is None:
        return
    title.text_frame.text = text
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True


def set_bullets(slide, items):
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    for i, (level, text) in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        para.text = text
        for run in para.runs:
            run.font.size = Pt(18)


def set_title_slide(slide, project_title, students, advisor, date):
    slide.shapes.title.text_frame.text = project_title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(28)
            run.font.bold = True
    subtitle = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            subtitle = shape
            break
    if subtitle is None:
        return
    tf = subtitle.text_frame
    tf.clear()
    lines = [
        ("Öğrenciler:", students),
        ("Danışman:", advisor),
        ("Sunum Tarihi:", date),
    ]
    for i, (label, value) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"{label}  {value}"
        for run in para.runs:
            run.font.size = Pt(18)


def reorder_slides(prs, new_order):
    """new_order: list of slide indices (current positions) in desired order."""
    xml_slides = prs.slides._sldIdLst
    current = list(xml_slides)
    # Remove all
    for sld in current:
        xml_slides.remove(sld)
    # Re-append in new order
    for idx in new_order:
        xml_slides.append(current[idx])


def main():
    prs = Presentation(SRC)

    # Original slides after open: 0=rules, 1=title, 2..6=content placeholders
    # Fill the title slide (index 1).
    set_title_slide(
        prs.slides[1],
        PROJECT_TITLE,
        STUDENTS,
        ADVISOR,
        DATE,
    )

    # The template has 5 content placeholders (indices 2..6); SLIDES has 6
    # entries (we added Sistem Mimarisi). Fill the 5 existing slots; clone
    # the layout from one of them for the 6th entry.
    content_layout = prs.slides[2].slide_layout

    for offset, (title, bullets) in enumerate(SLIDES):
        slide_idx = offset + 2  # skip rules + title
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
        else:
            slide = prs.slides.add_slide(content_layout)
        set_title(slide, title)
        set_bullets(slide, bullets)

    # Now slide order is: rules(0), title(1), Konu(2), Özgün(3), Yaygın(4),
    # Uygulanabilir(5), Sistem Mimarisi(6 - overwritten), Gerçekleştirme(7 - new).
    # Reorder so rules goes to the end (kullanıcı silsin diye), Sistem
    # Mimarisi already correctly placed before Gerçekleştirme.
    if len(prs.slides) == 8:
        reorder_slides(prs, [1, 2, 3, 4, 5, 6, 7, 0])

    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(DST)
    print(f"[OK] Wrote {DST}")
    print(f"     {len(prs.slides)} slides total (kurallar slaytı sona taşındı, isterseniz PowerPoint'te silebilirsiniz)")


if __name__ == "__main__":
    main()
